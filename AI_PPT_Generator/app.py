import os
import io
import re
import json
from socket import CAPI
import requests
from flask import Flask, render_template, request, send_file
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from io import BytesIO
from dotenv import load_dotenv
import os

# Load .env file
load_dotenv()

app = Flask(__name__, static_folder="static", template_folder="templates")

# Use env var for secret key
# API_KEY can be set via environment variable or hardcoded here (not recommended for production)
API_KEY = os.environ.get("SECRET_KEY", "default_secret_key")
API_KEY = os.environ.get("OPENROUTER_KEY", None)
OPENROUTER_URL = "https://openrouter.ai/api/v1/chat/completions"
LOGO_PATH = None  # optional: "static/logo.png"

# Theme palettes (matches CSS preview names)
THEMES = {
    "modern_blue": {
        "bg": RGBColor(245, 247, 250),
        "sidebar": RGBColor(11, 102, 182),
        "title": RGBColor(4, 58, 107),
        "bullet": RGBColor(33, 37, 41),
        "font": "Calibri"
    },
    "warm_retro": {
        "bg": RGBColor(255, 250, 245),
        "sidebar": RGBColor(232, 106, 42),
        "title": RGBColor(102, 51, 0),
        "bullet": RGBColor(64, 32, 0),
        "font": "Georgia"
    },
    "aqua_minimal": {
        "bg": RGBColor(245, 255, 255),
        "sidebar": RGBColor(30, 166, 168),
        "title": RGBColor(0, 77, 77),
        "bullet": RGBColor(0, 51, 51),
        "font": "Arial"
    },
    "dark_elegant": {
        "bg": RGBColor(20, 22, 25),
        "sidebar": RGBColor(68, 75, 255),
        "title": RGBColor(255, 255, 255),
        "bullet": RGBColor(220, 220, 220),
        "font": "Calibri"
    }
}

# -------------------------
# Utilities: image fetcher
# -------------------------
def safe_get_image(url):
    try:
        headers = {"User-Agent": "pptx-generator/1.0"}
        resp = requests.get(url, headers=headers, timeout=7, allow_redirects=True)
        ctype = resp.headers.get("Content-Type", "")
        if resp.status_code == 200 and ctype.startswith("image"):
            return BytesIO(resp.content)
    except Exception as e:
        print("safe_get_image error:", e)
    return None

# -------------------------
# Slide builder
# -------------------------
def add_slide(prs, title, bullets, palette, logo_path=None):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    # background
    try:
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = palette["bg"]
    except Exception:
        pass

    # left accent
    try:
        slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.6), prs.slide_height).fill.solid()
        left = slide.shapes[-1]
        left.fill.fore_color.rgb = palette["sidebar"]
        left.line.fill.background()
    except Exception:
        pass

    # title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(6.5), Inches(1.0))
    title_tf = title_box.text_frame
    title_tf.clear()
    p = title_tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.name = palette.get("font", "Calibri")
    p.font.color.rgb = palette["title"]
    p.alignment = PP_ALIGN.LEFT

    # bullets
    body_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(5.8), Inches(4.4))
    body_tf = body_box.text_frame
    body_tf.word_wrap = True
    body_tf.clear()
    for b in bullets:
        text = re.sub(r'^[\-\•\*\d\.\)\s]+', '', b).strip()
        if not text:
            continue
        para = body_tf.add_paragraph()
        para.text = text
        para.level = 0
        para.font.size = Pt(18)
        para.font.name = palette.get("font", "Calibri")
        para.font.color.rgb = palette["bullet"]

    # image on right (Unsplash)
    try:
        query = (title or "abstract").replace(" ", "%20")
        image_url = f"https://source.unsplash.com/800x600/?{query}"
        img_stream = safe_get_image(image_url)
        if img_stream:
            slide.shapes.add_picture(img_stream, Inches(6.2), Inches(1.2), width=Inches(3.0), height=Inches(2.2))
    except Exception as e:
        print("image insertion error:", e)

    # optional logo
    if logo_path and os.path.exists(logo_path):
        try:
            slide.shapes.add_picture(logo_path, prs.slide_width - Inches(1.2), prs.slide_height - Inches(0.9), width=Inches(0.9))
        except Exception as e:
            print("logo error:", e)

# -------------------------
# AI call and parsing
# -------------------------
def get_ai_slides(topic):
    if not API_KEY:
        return [("Error", ["No API key set. Set OPENROUTER_KEY environment variable."])]
    headers = {"Authorization": f"Bearer {API_KEY}", "Content-Type": "application/json"}
    prompt = f"""Create a PowerPoint presentation on the topic: \"{topic}\".
Structure it into 5 slides. Each slide should have a title and 3–5 bullet points.
Return the response like this:
Slide 1: Title
- Bullet
- Bullet
Slide 2: Title
- Bullet
..."""
    payload = {"model": "openai/gpt-3.5-turbo", "max_tokens": 1500, "messages": [{"role": "user", "content": prompt}]}

    try:
        res = requests.post(OPENROUTER_URL, headers=headers, data=json.dumps(payload), timeout=18)
    except Exception as e:
        return [("Error", [f"API request failed: {e}"])]

    if res.status_code != 200:
        return [("Error", [f"API error: {res.status_code}", res.text[:800]])]

    try:
        content = res.json()['choices'][0]['message']['content']
    except Exception as e:
        return [("Error", [f"Response parse error: {e}"])]

    slides = []
    cur_title, cur_bullets = None, []
    for raw in content.splitlines():
        line = raw.strip()
        if not line:
            continue
        m = re.match(r'^(?:Slide\s*\d+[:\-\)]\s*)(.+)$', line, re.IGNORECASE)
        if m:
            if cur_title:
                slides.append((cur_title, cur_bullets))
            cur_title = m.group(1).strip()
            cur_bullets = []
            continue
        if line.startswith('-') or line.startswith('•') or line.startswith('*'):
            cur_bullets.append(line)
            continue
        m2 = re.match(r'^\d+\.\s+(.+)$', line)
        if m2:
            if cur_title:
                slides.append((cur_title, cur_bullets))
            cur_title = m2.group(1).strip()
            cur_bullets = []
            continue
        if cur_title and len(cur_bullets) < 6 and len(line.split()) < 20:
            cur_bullets.append(line)
            continue
        if not cur_title:
            cur_title = line
            cur_bullets = []
        else:
            cur_bullets.append(line)
    if cur_title:
        slides.append((cur_title, cur_bullets))

    # ensure 5 slides
    if len(slides) < 5:
        while len(slides) < 5:
            slides.append(("Additional Info", ["More details coming soon...", "Edit this slide", "Customize content"]))
    elif len(slides) > 7:
        slides = slides[:7]

    return slides

# -------------------------
# Routes
# -------------------------
@app.route("/", methods=["GET"])
def index():
    return render_template("index.html")

@app.route("/select_theme", methods=["POST"])
def select_theme():
    topic = request.form.get("topic", "").strip()
    if not topic:
        return render_template("index.html", error="Please enter a topic.")
    return render_template("select_theme.html", topic=topic)

@app.route("/generate_ppt", methods=["POST"])
def generate_ppt():
    topic = request.form.get("topic", "Untitled Presentation").strip()
    theme_key = request.form.get("theme", "modern_blue")
    palette = THEMES.get(theme_key, THEMES["modern_blue"])

    slides = get_ai_slides(topic)

    prs = Presentation()
    for title, bullets in slides:
        add_slide(prs, title, bullets, palette=palette, logo_path=LOGO_PATH)

    ppt_io = BytesIO()
    prs.save(ppt_io)
    ppt_io.seek(0)
    safe_name = f"{topic.replace(' ', '_')}.pptx"
    return send_file(ppt_io, as_attachment=True, download_name=safe_name)

@app.route("/health")
def health():
    return "OK", 200

if __name__ == "__main__":
    app.run(debug=True, host="0.0.0.0", port=5000)
